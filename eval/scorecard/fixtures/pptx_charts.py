#!/usr/bin/env python3
"""Generate ``charts.pptx``: a title slide and three chart slides.

    uv run --with python-pptx python eval/scorecard/fixtures/pptx_charts.py [out_dir]

Slide 2 holds a two-series clustered column chart with a chart title and
axis titles, slide 3 a pie chart with a title, slide 4 a line chart with no
chart title. Every slide has a title placeholder, so the deck title and the
slide headings are the heading tree. Data and titles come from
``chart_data.py``. Deterministic: fixed core properties, and the package
plus each chart's embedded workbook are rewritten with fixed zip
timestamps and a fixed workbook creation date (python-pptx writes the
embedded workbook through xlsxwriter, which offers no clock hook).
"""

from __future__ import annotations

import sys
from datetime import datetime, timezone
from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.util import Inches

HERE = Path(__file__).resolve().parent
if str(HERE) not in sys.path:
    sys.path.insert(0, str(HERE))

from chart_data import BAR, LINE, PIE, ChartSpec, normalize_zip  # noqa: E402

DEFAULT_OUT = HERE.parents[2] / "tests" / "golden" / "corpus"
FIXED_TIME = datetime(2000, 1, 1, tzinfo=timezone.utc)
DECK_TITLE = "Charts Fixture Deck"
SUBTITLE = "Three chart slides: a column chart, a pie chart and a line chart"

# (slide title, chart, chart type, whether the chart carries its own title)
SLIDES: tuple[tuple[str, ChartSpec, XL_CHART_TYPE, bool], ...] = (
    ("Quarterly Revenue by Region", BAR, XL_CHART_TYPE.COLUMN_CLUSTERED, True),
    ("Market Share", PIE, XL_CHART_TYPE.PIE, True),
    ("Temperature Trend", LINE, XL_CHART_TYPE.LINE_MARKERS, False),
)
PIE_TITLE = "Share by segment"


def chart_title_for(spec: ChartSpec, titled: bool) -> str | None:
    if not titled:
        return None
    return spec.title or PIE_TITLE


def add_chart_slide(prs: Presentation, slide_title: str, spec: ChartSpec, kind: XL_CHART_TYPE,
                    titled: bool) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = slide_title
    data = CategoryChartData()
    data.categories = list(spec.categories)
    for index in range(spec.series_count):
        data.add_series(spec.header[index + 1], list(spec.series(index)))
    frame = slide.shapes.add_chart(kind, Inches(0.7), Inches(1.6), Inches(8.6), Inches(5.0), data)
    chart = frame.chart
    title = chart_title_for(spec, titled)
    if title:
        chart.has_title = True
        chart.chart_title.text_frame.text = title
    else:
        chart.has_title = False
    chart.has_legend = spec.series_count > 1 or spec.kind == "pie"
    if chart.has_legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
    if spec.kind != "pie":
        if spec.x_title:
            chart.category_axis.has_title = True
            chart.category_axis.axis_title.text_frame.text = spec.x_title
        if spec.y_title:
            chart.value_axis.has_title = True
            chart.value_axis.axis_title.text_frame.text = spec.y_title


def build() -> Presentation:
    prs = Presentation()
    title = prs.slides.add_slide(prs.slide_layouts[0])
    title.shapes.title.text = DECK_TITLE
    title.placeholders[1].text = SUBTITLE
    for slide_title, spec, kind, titled in SLIDES:
        add_chart_slide(prs, slide_title, spec, kind, titled)
    props = prs.core_properties
    props.author = ""
    props.last_modified_by = ""
    props.created = FIXED_TIME
    props.modified = FIXED_TIME
    props.title = DECK_TITLE
    return prs


def main(argv: list[str]) -> int:
    out_dir = Path(argv[0]) if argv else DEFAULT_OUT
    out_dir.mkdir(parents=True, exist_ok=True)
    target = out_dir / "charts.pptx"
    build().save(str(target))
    normalize_zip(target)
    print(f"{target} ({target.stat().st_size} bytes)")
    return 0


if __name__ == "__main__":
    sys.exit(main(sys.argv[1:]))
