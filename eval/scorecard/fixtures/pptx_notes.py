#!/usr/bin/env python3
"""Generate ``notes.pptx``: title, bullets, a table and speaker notes on three slides.

    uv run --with python-pptx python eval/scorecard/fixtures/pptx_notes.py [out_dir]

Deterministic: core properties carry a fixed timestamp and no author.
"""

from __future__ import annotations

import sys
from datetime import UTC, datetime
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

HERE = Path(__file__).resolve().parent
DEFAULT_OUT = HERE.parents[2] / "tests" / "golden" / "corpus"
FIXED_TIME = datetime(2000, 1, 1, tzinfo=UTC)

BULLETS = (
    (0, "Three kinds of slide content"),
    (1, "a title with a subtitle"),
    (1, "a bulleted list with two levels"),
    (1, "a four column table"),
    (0, "Speaker notes on every content slide"),
    (1, "notes are not on the slide canvas"),
)
TABLE = (
    ("Station", "Line", "Opened", "Platforms"),
    ("Harbor Point", "Blue", "1912", "2"),
    ("Mill Street", "Blue", "1925", "4"),
    ("Orchard Hill", "Green", "1958", "2"),
)
NOTES = {
    0: "Welcome. This deck exists so a parser can be checked for speaker notes; read this note aloud only in a test.",
    1: "Second slide notes: the two-level bullets are the thing to check, and the notes should not be mixed into them.",
    2: "Third slide notes: the table has one header row and three data rows; Orchard Hill opened in 1958.",
}


def build() -> Presentation:
    prs = Presentation()
    title = prs.slides.add_slide(prs.slide_layouts[0])
    title.shapes.title.text = "Notes Fixture Deck"
    title.placeholders[1].text = "A small deck with speaker notes, bullets and a table"

    bullets = prs.slides.add_slide(prs.slide_layouts[1])
    bullets.shapes.title.text = "What This Deck Covers"
    frame = bullets.placeholders[1].text_frame
    frame.text = BULLETS[0][1]
    for level, text in BULLETS[1:]:
        paragraph = frame.add_paragraph()
        paragraph.text = text
        paragraph.level = level

    table_slide = prs.slides.add_slide(prs.slide_layouts[5])
    table_slide.shapes.title.text = "Stations by Line"
    shape = table_slide.shapes.add_table(len(TABLE), len(TABLE[0]), Inches(0.7), Inches(1.8), Inches(8.5), Inches(2.0))
    for r, row in enumerate(TABLE):
        for c, value in enumerate(row):
            cell = shape.table.cell(r, c)
            cell.text = value
            cell.text_frame.paragraphs[0].font.size = Pt(16)

    closing = prs.slides.add_slide(prs.slide_layouts[5])
    closing.shapes.title.text = "Closing Slide Without Notes"

    for index, text in NOTES.items():
        prs.slides[index].notes_slide.notes_text_frame.text = text

    props = prs.core_properties
    props.author = ""
    props.last_modified_by = ""
    props.created = FIXED_TIME
    props.modified = FIXED_TIME
    props.title = "Notes Fixture Deck"
    return prs


def main(argv: list[str]) -> int:
    out_dir = Path(argv[0]) if argv else DEFAULT_OUT
    out_dir.mkdir(parents=True, exist_ok=True)
    target = out_dir / "notes.pptx"
    build().save(str(target))
    print(f"{target} ({target.stat().st_size} bytes)")
    return 0


if __name__ == "__main__":
    sys.exit(main(sys.argv[1:]))
